from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


ROOT = Path(__file__).resolve().parents[1]
PPTX_OUT = ROOT / "INVESTOR_PITCH_DECK.pptx"
PDF_OUT = ROOT / "INVESTOR_PITCH_SUMMARY.pdf"


BG = RGBColor(247, 248, 252)
NAVY = RGBColor(13, 36, 76)
TEAL = RGBColor(32, 164, 153)
GOLD = RGBColor(241, 180, 76)
INK = RGBColor(28, 34, 52)
MUTED = RGBColor(93, 100, 120)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(228, 236, 246)


def add_background(slide, accent=False):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    top_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.6)
    )
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = NAVY
    top_band.line.fill.background()

    if accent:
        band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.3), Inches(0.6), Inches(4.1), Inches(0.18)
        )
        band.fill.solid()
        band.fill.fore_color.rgb = TEAL
        band.line.fill.background()


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.9), Inches(8.8), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.72), Inches(9.5), Inches(0.5))
        p = sub_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = "Aptos"
        r.font.size = Pt(12)
        r.font.color.rgb = MUTED


def add_bullets(slide, bullets, left=0.9, top=2.0, width=7.0, height=4.5, level0_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True

    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"• {item}"
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(level0_size)
        p.font.color.rgb = INK
        p.space_after = Pt(8)


def add_quote_card(slide, title, body, x, y, w, h, fill_rgb, title_rgb=WHITE, body_rgb=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = fill_rgb

    tbox = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.18), Inches(w - 0.5), Inches(0.45))
    p = tbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = title_rgb

    bbox = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.62), Inches(w - 0.5), Inches(h - 0.8))
    p = bbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = body
    r.font.name = "Aptos"
    r.font.size = Pt(13)
    r.font.color.rgb = body_rgb


def add_metric_card(slide, value, label, x, y, w=2.5, h=1.5):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = PALE

    vbox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.18), Inches(w - 0.4), Inches(0.55))
    p = vbox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = NAVY

    lbox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.78), Inches(w - 0.4), Inches(0.4))
    p = lbox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.color.rgb = MUTED


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(6.5), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, accent=True)
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.1), Inches(6.8), Inches(4.9))
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.color.rgb = NAVY
    add_quote_card(
        slide,
        "NEXORA",
        "Building the assessment operating system for institutes, with a clear path to student progression, monetization, and long-term learning infrastructure.",
        0.95,
        1.45,
        6.2,
        2.2,
        NAVY,
    )
    add_quote_card(
        slide,
        "Stage Status",
        "Backend, frontend, SSL, Nginx, and role-aware workflows are already deployed and demoable.",
        0.95,
        3.95,
        3.0,
        1.25,
        TEAL,
    )
    add_quote_card(
        slide,
        "Current Wedge",
        "Institute-first assessment infrastructure before broad consumer expansion.",
        4.15,
        3.95,
        2.95,
        1.25,
        GOLD,
        title_rgb=NAVY,
        body_rgb=NAVY,
    )
    add_title(slide, "Investor Pitch", "Assessment infrastructure first. Student progression platform next.")
    add_metric_card(slide, "B2B First", "School / institute wedge", 8.0, 1.5)
    add_metric_card(slide, "Stage Live", "Deployed beta stack", 10.55, 1.5)
    add_metric_card(slide, "Modular", "Assessment + analytics + economy", 8.0, 3.3, 5.05)
    add_footer(slide, "Nexora | Investor conversation draft")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "The Problem", "Institutes need more than content. They need dependable assessment operations.")
    add_bullets(
        slide,
        [
            "Question creation, testing, evaluation, and analytics still live across fragmented tools.",
            "Teachers spend too much time on execution and too little on actual learning improvement.",
            "Students get delayed or shallow feedback, not a structured improvement loop.",
            "Most platforms either sell content, sell a generic LMS, or ignore exam execution quality.",
        ],
        width=7.4,
        height=4.8,
        level0_size=20,
    )
    add_quote_card(slide, "Gap in the market", "Operationally strong assessment infrastructure for institutes is still under-served.", 8.4, 2.1, 4.1, 2.3, TEAL)
    add_footer(slide, "Problem: digital assessment is still fragmented and weakly operationalized")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Our Solution", "Nexora is a role-aware assessment and progression backbone.")
    add_bullets(
        slide,
        [
            "Institute-scoped academic setup with teacher and student role control.",
            "Question bank, exam builder, publishing, attempts, results, and analytics in one stack.",
            "Student-facing workflows for attempts, review, results, analytics, and weak-area visibility.",
            "Configurable economy and unlock foundations for future monetization and premium layers.",
        ],
        width=7.0,
        height=4.8,
        level0_size=19,
    )
    add_metric_card(slide, "Teachers", "Create, assign, evaluate", 8.1, 2.0)
    add_metric_card(slide, "Students", "Attempt, review, improve", 10.65, 2.0)
    add_metric_card(slide, "Institutes", "Operate at scale", 8.1, 4.0)
    add_metric_card(slide, "Platform", "Control rules centrally", 10.65, 4.0)
    add_footer(slide, "Solution: one system for exam operations, feedback loops, and future monetization")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "What Is Already Built", "This is not just an idea. The product is demoable today.")
    add_bullets(
        slide,
        [
            "Django backend with institute scope, roles, exams, attempts, results, and analytics.",
            "Next.js web frontend for student and teacher workflows.",
            "Stage deployment completed with SSL, Nginx, systemd, PostgreSQL, and production runbooks.",
            "Seed flows for institutes, economy defaults, showcase questions, and exam demos.",
        ],
        width=7.2,
        height=4.8,
        level0_size=19,
    )
    add_metric_card(slide, "Teacher Web", "Dashboard, exams, bank, results", 8.2, 2.1)
    add_metric_card(slide, "Student Web", "Attempts, results, analytics", 10.8, 2.1)
    add_metric_card(slide, "Stage Ops", "Backend + frontend + SSL live", 8.2, 4.1, 5.1)
    add_footer(slide, "Current product status: stage-deployed and beta-ready")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Why This Wedge Works", "Assessment infrastructure creates real operational stickiness.")
    add_bullets(
        slide,
        [
            "Exam workflows are business-critical and recurring for every institute.",
            "Results and analytics create repeated teacher and student engagement.",
            "Once question banks, users, and academic structures are in the system, switching costs rise.",
            "This creates a stronger path to paid adoption than a content-only story.",
        ],
        width=7.0,
        height=4.8,
        level0_size=19,
    )
    add_quote_card(slide, "Positioning", "We are not another generic LMS. We are building the assessment operating system that institutes rely on day to day.", 8.0, 2.2, 4.5, 2.8, NAVY)
    add_footer(slide, "Operational wedge before broad edtech sprawl")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Business Model", "B2B first, with a credible path to hybrid monetization.")
    add_quote_card(slide, "Phase 1", "Institute subscription for assessment operations", 0.8, 2.0, 3.8, 1.6, TEAL)
    add_quote_card(slide, "Phase 2", "Premium exam bundles, analytics, and readiness products", 4.8, 2.0, 3.8, 1.6, NAVY)
    add_quote_card(slide, "Phase 3", "Student unlocks, rewards, and subscription-ready economy layers", 8.8, 2.0, 3.8, 1.6, GOLD, title_rgb=NAVY, body_rgb=NAVY)
    add_bullets(
        slide,
        [
            "Core monetization starts with institute utility, not speculative consumer virality.",
            "The backend is already being shaped for configurable stars, unlocks, entitlements, and plans.",
            "This allows future monetization without rebuilding the product foundation.",
        ],
        left=0.95,
        top=4.3,
        width=11.2,
        height=2.0,
        level0_size=18,
    )
    add_footer(slide, "Commercial path: institute SaaS first, student monetization second")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Market Entry", "Start with institutes. Expand through operational trust.")
    add_bullets(
        slide,
        [
            "Pilot with schools, coaching centers, and institute operators who already run regular assessments.",
            "Use teacher and student workflows to validate recurring usage and institutional fit.",
            "Convert operational adoption into paid SaaS relationships before broader public expansion.",
            "Layer premium readiness, practice, and progression journeys once the core wedge is sticky.",
        ],
        width=7.0,
        height=4.6,
        level0_size=19,
    )
    add_metric_card(slide, "Pilot", "Design-partner institutions", 8.2, 2.2)
    add_metric_card(slide, "Validate", "Usage + workflow fit", 10.75, 2.2)
    add_metric_card(slide, "Expand", "Student progression layers", 8.2, 4.2, 5.05)
    add_footer(slide, "Go-to-market: founder-led pilot execution, then repeatable institutional sales")

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Roadmap", "Clear next milestones without pretending the whole vision is already complete.")
    add_bullets(
        slide,
        [
            "Harden platform admin and institute admin product surfaces.",
            "Expand curated question libraries and showcase exam depth.",
            "Run beta validation with real institutes on the live stage environment.",
            "Activate economy-backed premium readiness and unlock flows in the next phase.",
        ],
        width=7.2,
        height=4.8,
        level0_size=19,
    )
    add_quote_card(slide, "Not the focus right now", "ERP, payroll, transport, live classes, broad consumer onboarding, and AI tutoring are intentionally out of current launch scope.", 8.0, 2.3, 4.5, 2.7, WHITE, title_rgb=NAVY, body_rgb=INK)
    add_footer(slide, "Focused roadmap: strengthen the backbone, then expand deliberately")

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "The Ask", "Fund product hardening, pilot execution, and early institutional traction.")
    add_bullets(
        slide,
        [
            "Product engineering and reliability hardening.",
            "UI and workflow polish for admin, teacher, and student surfaces.",
            "Pilot onboarding and implementation support.",
            "Assessment content operations and customer discovery.",
        ],
        width=6.8,
        height=4.2,
        level0_size=20,
    )
    add_quote_card(slide, "Investor takeaway", "Nexora is a serious academic infrastructure product with a clear operational wedge and a credible path to student-scale monetization.", 8.0, 2.2, 4.6, 2.7, NAVY)
    add_metric_card(slide, "Now", "Beta-ready stage product", 8.0, 5.2)
    add_metric_card(slide, "Next", "Institutional validation", 10.55, 5.2)
    add_footer(slide, "Ask: help convert a working backbone into a commercially validated product")

    prs.save(PPTX_OUT)


def build_pdf():
    doc = SimpleDocTemplate(str(PDF_OUT), pagesize=A4, rightMargin=42, leftMargin=42, topMargin=42, bottomMargin=36)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="TitleNexora", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=24, textColor=colors.HexColor("#0D244C"), spaceAfter=18))
    styles.add(ParagraphStyle(name="HeadingNexora", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=14, textColor=colors.HexColor("#0D244C"), spaceBefore=14, spaceAfter=8))
    styles.add(ParagraphStyle(name="BodyNexora", parent=styles["BodyText"], fontName="Helvetica", fontSize=10.5, leading=15, textColor=colors.HexColor("#1C2234")))

    story = []
    story.append(Paragraph("Nexora Investor Summary", styles["TitleNexora"]))
    story.append(Paragraph("Assessment infrastructure first. Student progression platform next.", styles["BodyNexora"]))
    story.append(Spacer(1, 0.18 * inch))

    sections = [
        ("What Nexora Is", "Nexora is an institute-first assessment, exam-readiness, and academic progression platform. It helps institutes run question banks, exams, attempts, results, and analytics in one operating system."),
        ("The Core Problem", "Institutes still manage assessments through fragmented tools, delayed evaluation, weak feedback loops, and manual academic operations. Students do not receive structured improvement visibility, and institutions cannot easily scale assessment into a digital product."),
        ("Our Solution", "Nexora provides a role-aware platform for platform admins, institute admins, teachers, and students. The current product already supports institute-scoped academic setup, exam creation, attempt workflows, result generation, analytics, and teacher operations."),
        ("Why This Matters", "Most edtech products either focus on static content or generic LMS workflows. Nexora focuses on the operational assessment layer, which is more recurring, more institutionally sticky, and a stronger wedge for monetization."),
        ("Current Product Status", "The backend and frontend are working, the platform is stage-deployed, and the system already supports production-style deployment with Nginx, SSL, systemd, PostgreSQL, and role-aware product surfaces."),
        ("Business Model", "The near-term business model is institute subscription revenue. The longer-term upside includes premium student readiness products, unlocks, reward systems, and configurable monetization on top of the existing assessment backbone."),
        ("Go-To-Market", "Start with pilot schools, coaching centers, and institutes that already run regular assessments. Prove recurring usage through teacher and student workflows, then expand into paid institutional adoption and premium student features."),
        ("Why We Can Win", "Nexora already has a meaningful backend foundation, a clear product wedge, and an architecture designed for future monetization without requiring a backend rewrite."),
        ("What We Need Next", "The next phase is product hardening, pilot execution, admin workflow polish, and institutional validation. The goal is to turn a strong operational backbone into a commercially repeatable product."),
    ]

    for heading, body in sections:
        story.append(Paragraph(heading, styles["HeadingNexora"]))
        story.append(Paragraph(body, styles["BodyNexora"]))

    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph("Current Product Snapshot", styles["HeadingNexora"]))

    table = Table(
        [
            ["Area", "Status"],
            ["Backend", "Django exam, attempt, result, analytics, and role system working"],
            ["Frontend", "Next.js student and teacher product surfaces working"],
            ["Deployment", "Stage server live with HTTPS and service automation"],
            ["Monetization foundation", "Economy and unlock architecture prepared for next phase"],
        ],
        colWidths=[1.8 * inch, 4.9 * inch],
    )
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0D244C")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 9.5),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#F7F8FC")),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#DDE5F0")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 8),
                ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]
        )
    )
    story.append(table)
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph("Prepared from the live Nexora product build and deployment runbooks.", styles["BodyNexora"]))

    doc.build(story)


if __name__ == "__main__":
    build_pptx()
    build_pdf()
    print(f"Created: {PPTX_OUT}")
    print(f"Created: {PDF_OUT}")
